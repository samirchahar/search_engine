# src/extractor/extractor.py
# Handles text extraction from PDF, TXT, and DOCX files.
# PDF: split by page (natural boundary)
# TXT and DOCX: split into 500-word chunks with 100-word overlap
# Overlap prevents important sentences from being lost at chunk boundaries.

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image
import io

# Tell pytesseract where Tesseract is installed
pytesseract.pytesseract.tesseract_cmd = r'D:\College\TesseractOCR\tesseract.exe'


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> list[dict]:
    """
    Split a long text into overlapping chunks.
    Each chunk is chunk_size words, with overlap words repeated
    from the end of the previous chunk at the start of the next.
    Returns: [{"page": 1, "text": "..."}, ...]
    """
    words = text.split()
    if not words:
        return []

    results = []
    start = 0
    page_num = 1

    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = words[start:end]
        results.append({
            "page": page_num,
            "text": " ".join(chunk)
        })
        # Move forward by chunk_size minus overlap
        # This means the next chunk starts 100 words before this one ended
        start += chunk_size - overlap
        page_num += 1

    return results


def extract_pdf(filepath: str) -> list[dict]:
    """
    Extract text from a PDF file, one entry per page.
    Automatically falls back to OCR if a page has no extractable text.
    Returns: [{"page": 1, "text": "..."}, ...]
    """
    results = []
    doc = fitz.open(filepath)
    needs_ocr = []

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text()
        if text.strip():
            results.append({
                "page": page_num,
                "text": text
            })
        else:
            needs_ocr.append(page_num)
    doc.close()

    if needs_ocr:
        print(f"  Running OCR on {len(needs_ocr)} scanned page(s)...")
        ocr_results = extract_pdf_ocr(filepath)
        # Only add OCR results for pages that had no text
        for r in ocr_results:
            if r["page"] in needs_ocr:
                results.append(r)
        # Re-sort by page number
        results.sort(key=lambda x: x["page"])

    return results


def extract_pdf_ocr(filepath: str) -> list[dict]:
    """
    OCR fallback for scanned PDFs where text extraction returns nothing.
    Renders each page as an image and runs Tesseract OCR on it.
    Returns: [{"page": 1, "text": "..."}, ...]
    """
    results = []
    doc = fitz.open(filepath)
    for page_num, page in enumerate(doc, start=1):
        # Render page as image at 300 DPI for good OCR accuracy
        mat = fitz.Matrix(300/72, 300/72)
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.tobytes("png")
        image = Image.open(io.BytesIO(img_data))
        text = pytesseract.image_to_string(image)
        if text.strip():
            results.append({
                "page": page_num,
                "text": text
            })
    doc.close()
    return results


def extract_txt(filepath: str) -> list[dict]:
    """
    Extract text from a TXT file.
    Splits into 500-word chunks with 100-word overlap.
    Returns: [{"page": 1, "text": "..."}, ...]
    """
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    if not text.strip():
        return []
    return chunk_text(text)


def extract_docx(filepath: str) -> list[dict]:
    """
    Extract text from a DOCX file.
    Splits into 500-word chunks with 100-word overlap.
    Page numbers are approximate — use the snippet to locate content in Word.
    Returns: [{"page": 1, "text": "..."}, ...]
    """
    doc = Document(filepath)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    full_text = " ".join(paragraphs)
    return chunk_text(full_text)


def extract_pptx(filepath: str) -> list[dict]:
    """
    Extract text from a PPTX file.
    Each slide is treated as one page.
    Returns: [{"page": 1, "text": "..."}, ...]
    """
    prs = Presentation(filepath)
    results = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            results.append({
                "page": slide_num,
                "text": " ".join(texts)
            })

    return results


def extract_file(filepath: str) -> list[dict]:
    """
    Route a file to the correct extractor based on extension.
    Supported: .pdf, .txt, .docx, .pptx
    Returns: [{"page": N, "text": "..."}, ...]
    """
    if filepath.lower().endswith(".pdf"):
        return extract_pdf(filepath)
    elif filepath.lower().endswith(".txt"):
        return extract_txt(filepath)
    elif filepath.lower().endswith(".docx"):
        return extract_docx(filepath)
    elif filepath.lower().endswith(".pptx"):
        return extract_pptx(filepath)
    else:
        print(f"Unsupported file type: {filepath}")
        return []